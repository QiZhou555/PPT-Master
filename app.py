import sys
import os
import configparser
from flask import Flask, render_template, request, jsonify, send_file, url_for, abort, make_response
import json
import time
import traceback
import random
from datetime import datetime
import logging

# 导入用户提供的文生文API模块
from openai import OpenAI

# 导入图片生成模块
from image_generator import generate_image, generate_image_prompts

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# 调试信息
print(f"当前Python路径: {sys.executable}")
print(f"Python版本: {sys.version}")
print(f"模块搜索路径: {sys.path}")

# 读取配置信息
def load_config():
    """从环境变量或配置文件加载配置"""
    config = {
        'server': {
            'host': os.environ.get('FLASK_HOST', '0.0.0.0'),
            'port': int(os.environ.get('FLASK_PORT', 5001)),
            'secret_key': os.environ.get('FLASK_SECRET_KEY', 'default-secret-key')
        },
        'api': {
            'base_url': os.environ.get('API_BASE_URL', 'https://api-inference.modelscope.cn/v1'),
            'api_key': os.environ.get('API_KEY', 'ms-4b457ae8-4cfd-4504-8ec2-8dc2fb930454')
        }
    }
    
    # 尝试从配置文件加载
    if os.path.exists('config.ini'):
        try:
            parser = configparser.ConfigParser()
            parser.read('config.ini', encoding='utf-8')
            
            if 'server' in parser:
                if parser.has_option('server', 'host'):
                    config['server']['host'] = parser.get('server', 'host')
                if parser.has_option('server', 'port'):
                    config['server']['port'] = parser.getint('server', 'port')
                if parser.has_option('server', 'secret_key'):
                    config['server']['secret_key'] = parser.get('server', 'secret_key')
            
            if 'api' in parser:
                if parser.has_option('api', 'base_url'):
                    config['api']['base_url'] = parser.get('api', 'base_url')
                if parser.has_option('api', 'api_key'):
                    config['api']['api_key'] = parser.get('api', 'api_key')
                    
        except Exception as e:
            logger.warning(f"读取配置文件失败: {e}")
    
    return config

# 加载配置
config = load_config()

# 创建OpenAI客户端
client = OpenAI(
    base_url=config['api']['base_url'],
    api_key=config['api']['api_key'],
)

# 尝试导入必要的模块
PPTX_AVAILABLE = False
PDF_AVAILABLE = False
reportlab = None
colors = None
SimpleDocTemplate = None
Paragraph = None
Spacer = None
getSampleStyleSheet = None
ParagraphStyle = None
inch = None
letter = None
landscape = None

# 先导入可能存在的模块
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
    print("python-pptx模块导入成功")
except ImportError:
    print("无法导入python-pptx模块，将尝试生成PDF文件")
    print(f"导入错误: {traceback.format_exc()}")

# 创建Flask应用
app = Flask(__name__)

# 添加对OPTIONS请求的支持
@app.route('/generate_content', methods=['OPTIONS'])
@app.route('/create_ppt', methods=['OPTIONS'])
@app.route('/download/<filename>', methods=['OPTIONS'])
def handle_options(filename=None):
    response = make_response()
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type, Authorization'
    response.headers['Access-Control-Max-Age'] = '86400'  # 24小时
    return response

# 添加CORS支持的中间件，确保所有响应都有CORS头
@app.after_request
def add_cors_headers(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type, Authorization'
    return response

# 设置应用密钥
app.secret_key = config['server']['secret_key']

# 配置文件存储文件夹 - 使用ppt_files而不是uploads
app.config['UPLOAD_FOLDER'] = os.path.join(os.getcwd(), 'ppt_files')
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# 主题颜色映射
COLOR_MAP = {
    "blue": {"hex": "#0078D4", "rgb": (0, 120, 212)},  
    "red": {"hex": "#FF5733", "rgb": (255, 87, 51)},
    "green": {"hex": "#2ECC71", "rgb": (46, 204, 113)},
    "紫色": {"hex": "#4c0099", "rgb": (76, 0, 153)},
    "橙色": {"hex": "#ff8000", "rgb": (255, 128, 0)},
    "蓝色": {"hex": "#003366", "rgb": (0, 51, 102)}
}

# 将16进制颜色转换为RGB
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# 生成文本文件作为替代方案
def create_text_file(content, filename):
    """
    创建文本文件作为PPT/PDF的替代方案
    """
    text_filename = f'{filename}.txt'
    text_path = os.path.join(app.config['UPLOAD_FOLDER'], text_filename)
    
    with open(text_path, 'w', encoding='utf-8') as f:
        # 写入标题页
        f.write("=" * 80 + '\n')
        f.write(f"{' ' * 20}{content['title']}\n")
        f.write(f"{' ' * 25}{content.get('subtitle', '')}\n")
        f.write("=" * 80 + '\n\n')
        
        # 写入目录
        f.write("目录\n")
        f.write("-" * 80 + '\n')
        for i, section in enumerate(content.get('sections', []), 1):
            f.write(f"{i}. {section['title']}\n")
        f.write("\n" * 2)
        
        # 写入各章节内容
        for section in content.get('sections', []):
            f.write(section['title'] + '\n')
            f.write("-" * 80 + '\n')
            f.write(section['content'] + '\n')
            f.write("\n" * 2)
    
    logger.info(f"成功创建文本文件: {text_path}")
    return text_filename

# 生成PPT内容的函数
def generate_ppt_content(topic, style='business', length='medium'):
    """
    根据主题和风格调用文生文API生成PPT内容
    """
    logger.info(f"生成PPT内容: {topic}, 风格: {style}, 长度: {length}")
    
    # 设置章节数量根据长度
    if length == 'short':
        section_count = 3
        detail_level = "简要"
    elif length == 'long':
        section_count = 6
        detail_level = "详细"
    else:  # medium
        section_count = 4
        detail_level = "中等"
    
    # 构建提示词，要求生成结构化的PPT内容
    system_prompt = f"""你是一个专业的PPT内容生成助手。请根据用户提供的主题，生成一个{detail_level}详细程度的{style}风格PPT内容。
    请严格按照以下JSON格式输出，不要包含任何额外的文字说明：
    {{
        "title": "PPT标题",
        "subtitle": "PPT副标题",
        "sections": [
            {{
                "title": "章节标题1",
                "content": ["要点1", "要点2", "要点3"],
                "detail": "章节详细描述"
            }},
            {{
                "title": "章节标题2",
                "content": ["要点1", "要点2", "要点3"],
                "detail": "章节详细描述"
            }}
        ]
    }}
    请确保生成{section_count}个章节，每个章节包含标题、要点列表和详细描述。
    内容必须与主题高度相关，专业且有深度。"""
    
    user_prompt = f"主题: {topic}"
    
    try:
        # 调用文生文API生成内容
        logger.info(f"调用文生文API生成PPT内容")
        response = client.chat.completions.create(
            model='Qwen/Qwen3-Next-80B-A3B-Instruct',
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': user_prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        # 获取API返回的内容
        api_response = response.choices[0].message.content.strip()
        logger.info(f"API返回内容: {api_response[:100]}...")
        
        # 解析JSON响应
        content_data = json.loads(api_response)
        
        # 构建完整的PPT内容结构，包含图片提示词
        content = {
            'title': content_data['title'],
            'subtitle': content_data.get('subtitle', f'{style}风格演示文稿'),
            'sections': [],
            'title_image_prompt': f'{content_data["title"]}的专业{style}风格图片，高质量，清晰，专业',
            'image_prompts': [
                {'prompt': f'{content_data["title"]}的专业{style}风格图片，高质量，清晰，专业', 'type': 'title', 'section_index': 0}
            ]
        }
        
        # 处理每个章节，添加图片提示词
        for i, section in enumerate(content_data['sections']):
            # 确保content字段是列表格式
            section_content = section['content'] if isinstance(section['content'], list) else [section['content']]
            
            # 添加完整的章节信息
            full_section = {
                'title': section['title'],
                'content': section_content,
                'detail': section.get('detail', ''),
                'image_prompt': f'{content_data["title"]} {section["title"]} 相关专业图片，{style}风格，高质量，清晰'
            }
            content['sections'].append(full_section)
            
            # 添加章节图片提示词
            content['image_prompts'].append({
                'prompt': f'{content_data["title"]} {section["title"]} 相关专业图片，{style}风格，高质量，清晰',
                'type': 'section',
                'section_index': i
            })
        
        logger.info(f"成功通过API生成PPT内容，共{len(content['sections'])}个章节")
        return content
        
    except json.JSONDecodeError as e:
        logger.error(f"解析API返回的JSON失败: {e}")
        logger.error(f"API返回的原始内容: {api_response}")
        # 返回备用内容
        return generate_fallback_content(topic, style, length)
    except Exception as e:
        logger.error(f"调用文生文API时出错: {e}")
        logger.error(f"错误详情: {traceback.format_exc()}")
        # 返回备用内容
        return generate_fallback_content(topic, style, length)

# 备用的内容生成函数，当API调用失败时使用
def generate_fallback_content(topic, style='business', length='medium'):
    """
    当API调用失败时使用的备用内容生成函数
    """
    logger.warning("使用备用内容生成逻辑")
    
    # 基础内容结构
    content = {
        'title': topic,
        'subtitle': f'专业{style}风格PPT',
        'sections': [],
        'title_image_prompt': f'{topic}的专业{style}风格图片，高质量，清晰，专业',
        'image_prompts': [
            {'prompt': f'{topic}的专业{style}风格图片，高质量，清晰，专业', 'type': 'title', 'section_index': 0}
        ]
    }
    
    # 根据PPT长度确定章节数量
    if length == 'short':
        num_sections = 3
    elif length == 'long':
        num_sections = 6
    else:  # medium
        num_sections = 4
    
    # 简单的通用章节
    base_sections = [
        {'title': '概述', 'content': [f'{topic}简介', '主要目标', '关键要点']},
        {'title': '详细内容', 'content': [f'{topic}核心内容', '重要信息', '相关数据']},
        {'title': '应用场景', 'content': [f'{topic}应用领域', '实际案例', '使用方法']},
        {'title': '总结', 'content': [f'{topic}要点回顾', '未来展望', '下一步行动']},
        {'title': '附录', 'content': ['参考资料', '补充信息', '联系方式']},
        {'title': 'Q&A', 'content': ['常见问题', '讨论环节', '反馈收集']}
    ]
    
    # 添加章节内容
    for i, section_template in enumerate(base_sections[:num_sections]):
        section = {
            'title': section_template['title'],
            'content': section_template['content'],
            'detail': f'本章节介绍{topic}中的{section_template["title"]}相关内容。',
            'image_prompt': f'{topic} {section_template["title"]} 相关图片'
        }
        content['sections'].append(section)
        
        # 添加章节图片提示词
        content['image_prompts'].append({
            'prompt': f'{topic} {section_template["title"]} 相关图片',
            'type': 'section',
            'section_index': i
        })
    
    return content

@app.route('/')
def index():
    try:
        return render_template('index.html')
    except Exception as e:
        logger.error(f"渲染首页时出错: {e}")
        return "服务器内部错误", 500

@app.route('/generate_content', methods=['POST'])
def generate_content():
    try:
        # 增加详细的请求日志
        logger.info(f"接收到POST请求: {request.get_data()[:200]}...")  # 只记录前200个字符避免日志过长
        
        # 增加对请求体的验证
        if not request.is_json:
            logger.error("请求体不是JSON格式")
            return jsonify({'error': '请求体必须是JSON格式'}), 400
        
        # 打印完整的JSON数据（如果可用）
        try:
            json_data = request.get_json()
            logger.info(f"JSON数据: {json_data}")
        except Exception as e:
            logger.error(f"解析JSON时出错: {e}")
            return jsonify({'error': 'JSON格式无效'}), 400
        
        topic = request.json.get('topic', '')
        style = request.json.get('style', 'business')  # 获取风格参数
        length = request.json.get('length', 'medium')  # 获取长度参数
        
        logger.info(f"提取的参数 - 主题: '{topic}', 风格: '{style}', 长度: '{length}'")
        
        if not topic:
            logger.warning("主题为空或未提供")
            return jsonify({'error': '请输入PPT主题'}), 400
        
        logger.info(f"开始生成内容，主题: {topic}，风格: {style}，长度: {length}")
        content = generate_ppt_content(topic, style, length)
        logger.info(f"内容生成成功，返回数据长度: {len(str(content))}")
        return jsonify({'content': content})
    except Exception as e:
        logger.error(f"处理生成内容请求时出错: {e}")
        logger.error(f"错误详情: {traceback.format_exc()}")
        return jsonify({'error': f'处理请求时出错: {str(e)}'}), 500

@app.route('/create_ppt', methods=['POST'])
def create_ppt_route():
    try:
        data = request.json
        content = data.get('content')
        theme_color = data.get('theme_color', '蓝色')

        # 验证内容
        if not content or not isinstance(content, dict):
            return jsonify({'error': '无效的PPT内容'}), 400

        # 动态导入reportlab模块（仅在需要时）
        global PDF_AVAILABLE, colors, SimpleDocTemplate, Paragraph, Spacer, getSampleStyleSheet, ParagraphStyle, inch, letter, landscape
        if not PDF_AVAILABLE:
            try:
                from reportlab.lib.pagesizes import letter, landscape
                from reportlab.lib import colors
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import inch
                PDF_AVAILABLE = True
                print("reportlab模块导入成功")
            except ImportError:
                print("无法导入reportlab模块")

        # 生成文件名（使用时间戳避免冲突）
        timestamp = str(int(time.time()))
        filename = f"{content['title'].replace(' ', '_')}_{timestamp}"
        
        # 尝试按照优先级生成文件
        if PPTX_AVAILABLE:
            # 创建PPT文件
            ppt_filename = f'{filename}.pptx'
            ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_filename)
            
            # 创建演示文稿对象
            prs = Presentation()
            
            # 获取主题颜色
            color_info = COLOR_MAP.get(theme_color, COLOR_MAP["蓝色"])
            rgb_color = RGBColor(*color_info["rgb"])
            
            # 生成图片目录
            images_dir = os.path.join(app.config['UPLOAD_FOLDER'], f"images_{timestamp}")
            os.makedirs(images_dir, exist_ok=True)
            
            # 存储生成的图片路径
            image_paths = []
            
            # 为每个提示词生成图片
            if 'image_prompts' in content:
                for image_prompt in content['image_prompts']:
                    prompt = image_prompt.get('prompt', '')
                    section_type = image_prompt.get('type', 'section')
                    section_index = image_prompt.get('section_index', 0)
                    
                    # 生成图片文件名
                    if section_type == 'title':
                        image_filename = f"title_{timestamp}.jpg"
                    else:
                        image_filename = f"section_{section_index}_{timestamp}.jpg"
                    
                    image_path = os.path.join(images_dir, image_filename)
                    
                    # 生成图片
                    try:
                        success = generate_image(prompt, image_path)
                        if success:
                            image_paths.append({
                                'path': image_path,
                                'type': section_type,
                                'section_index': section_index
                            })
                            logger.info(f"成功生成图片: {image_path}")
                    except Exception as e:
                        logger.error(f"生成图片时出错: {e}")
            
            # 创建标题页
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = content['title']
            subtitle.text = content.get('subtitle', '')
            
            # 设置标题文本颜色
            title_text_frame = title.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = rgb_color
            
            # 尝试添加标题图片
            title_images = [img for img in image_paths if img['type'] == 'title']
            if title_images:
                img_path = title_images[0]['path']
                left = Inches(1)
                top = Inches(3)
                width = Inches(7)
                height = Inches(4)
                slide.shapes.add_picture(img_path, left, top, width, height)
            
            # 创建目录页
            table_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(table_slide_layout)
            title = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title.text = "目录"
            
            # 添加目录内容
            tf = content_shape.text_frame
            tf.clear()
            
            for i, section in enumerate(content.get('sections', []), 1):
                p = tf.add_paragraph()
                p.text = f"{i}. {section['title']}"
                p.level = 0
            
            # 为每个章节创建幻灯片，增加数据验证
            for i, section in enumerate(content.get('sections', [])):
                try:
                    content_slide_layout = prs.slide_layouts[3]  # 使用带有图片占位符的布局
                    slide = prs.slides.add_slide(content_slide_layout)
                    title = slide.shapes.title
                    
                    # 确保标题存在且为字符串
                    title.text = str(section.get('title', f'章节 {i+1}'))
                    
                    # 设置标题文本颜色
                    title_text_frame = title.text_frame
                    for paragraph in title_text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = rgb_color
                    
                    # 添加章节内容，增加数据类型检查
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.clear()
                    
                    # 添加章节要点，确保content是列表
                    section_content = section.get('content', [])
                    if isinstance(section_content, list):
                        for point in section_content:
                            p = tf.add_paragraph()
                            p.text = str(point)  # 确保文本内容是字符串
                            p.level = 1
                    
                    # 添加章节详细描述
                    if 'detail' in section and isinstance(section['detail'], str):
                        p = tf.add_paragraph()
                        p.text = f"\n{section['detail']}"
                        p.level = 0
                    
                    # 尝试添加章节图片
                    section_images = [img for img in image_paths if img['type'] == 'section' and img['section_index'] == i]
                    if section_images and len(slide.placeholders) > 2:
                        img_path = section_images[0]['path']
                        # 获取占位符的位置和大小
                        placeholder = slide.placeholders[2]
                        left = placeholder.left
                        top = placeholder.top
                        width = placeholder.width
                        height = placeholder.height
                        # 使用slide.shapes.add_picture添加图片
                        slide.shapes.add_picture(img_path, left, top, width, height)
                except Exception as e:
                    logger.error(f"创建章节幻灯片时出错 (章节 {i+1}): {e}")
                    # 即使单个章节出错，也继续创建其他章节
                    continue
            
            # 创建结束页
            end_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(end_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "谢谢观看"
            subtitle.text = f"{content['title']} - 专业分析报告"
            
            # 设置标题文本颜色
            title_text_frame = title.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = rgb_color
            
            # 保存PPT文件
            prs.save(ppt_path)
            logger.info(f"成功创建PPT文件: {ppt_path}")
            
            # 生成前端预览所需的幻灯片数据
            slides_preview = [
                {'title': content['title'], 'subtitle': content.get('subtitle', '')}
            ]
            
            # 添加目录预览
            slides_preview.append({'title': '目录', 'content': '\n'.join([f"{i+1}. {section['title']}" for i, section in enumerate(content.get('sections', []))])})
            
            # 添加章节预览
            for section in content.get('sections', []):
                section_content = ''
                # 添加要点
                if isinstance(section.get('content'), list):
                    for point in section.get('content', []):
                        section_content += f"• {point}\n"
                # 添加详细描述
                if isinstance(section.get('detail'), str):
                    section_content += f"\n{section['detail']}"
                
                slides_preview.append({
                    'title': section.get('title', '章节标题'),
                    'content': section_content
                })
            
            # 添加结束页预览
            slides_preview.append({'title': '谢谢观看', 'subtitle': f"{content['title']} - 专业分析报告"})
            
            return jsonify({
                'ppt_path': ppt_filename,
                'file_url': f'/download/{ppt_filename}',
                'pdf_available': PDF_AVAILABLE,
                'file_type': 'pptx',
                'slides': slides_preview
            })
        elif PDF_AVAILABLE:
            # 创建PDF文件
            pdf_filename = f'{filename}.pdf'
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
            
            # 创建PDF文档
            doc = SimpleDocTemplate(pdf_path, pagesize=landscape(letter))
            story = []
            styles = getSampleStyleSheet()
            
            # 获取主题颜色
            color_info = COLOR_MAP.get(theme_color, COLOR_MAP["蓝色"])
            
            # 创建自定义样式
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontSize=36,
                textColor=colors.HexColor(color_info["hex"]),
                spaceAfter=30
            )
            
            # 添加标题页
            story.append(Paragraph(content['title'], title_style))
            if 'subtitle' in content:
                story.append(Paragraph(content['subtitle'], styles['Heading2']))
            
            # 添加章节内容
            for section in content.get('sections', []):
                story.append(Spacer(1, 0.5*inch))
                story.append(Paragraph(section.get('title', '章节标题'), styles['Heading1']))
                # 确保content是字符串
                section_content = section.get('content', '')
                if isinstance(section_content, list):
                    section_content = '\n'.join(section_content)
                story.append(Paragraph(str(section_content), styles['BodyText']))
                
                # 添加详细描述
                if 'detail' in section and isinstance(section['detail'], str):
                    story.append(Paragraph(section['detail'], styles['BodyText']))
            
            # 构建PDF文档
            doc.build(story)
            logger.info(f"成功创建PDF文件: {pdf_path}")
            
            return jsonify({
                'ppt_path': pdf_filename,
                'file_url': f'/download/{pdf_filename}',
                'pdf_available': True,
                'file_type': 'pdf'
            })
        else:
            # 当没有PPTX和PDF支持时，创建文本文件作为替代方案
            text_filename = create_text_file(content, filename)
            
            return jsonify({
                'ppt_path': text_filename,
                'file_url': f'/download/{text_filename}',
                'pdf_available': False,
                'file_type': 'txt',
                'warning': '服务器暂不支持生成PPT或PDF文件，已生成文本版本供参考。请安装相应依赖包以获得更好的体验。'
            })
            
    except Exception as e:
        # 记录详细错误信息
        print(f"创建演示文稿时发生错误: {str(e)}")
        print(f"错误详情: {traceback.format_exc()}")
        
        # 即使出错，也尝试创建文本文件作为兜底方案
        try:
            text_filename = create_text_file(content, filename)
            return jsonify({
                'ppt_path': text_filename,
                'file_url': f'/download/{text_filename}',
                'pdf_available': False,
                'file_type': 'txt',
                'warning': f'生成高级格式失败: {str(e)}，已生成文本版本供参考。',
                'error': str(e)
            })
        except Exception as fallback_error:
            # 如果连文本文件都创建失败，返回友好的错误消息
            return jsonify({'error': f'创建演示文稿时发生错误: {str(e)}'}), 500

@app.route('/download/<filename>')
def download_file(filename):
    try:
        # 安全检查：确保文件名中不包含路径遍历字符
        if '..' in filename or os.path.isabs(filename):
            abort(400)
        
        # 构建完整的文件路径
        safe_filename = os.path.basename(filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            abort(404)
        
        # 检查文件扩展名是否合法 - 添加.txt支持
        if not (file_path.endswith('.pptx') or file_path.endswith('.pdf') or file_path.endswith('.txt')):
            abort(400)
        
        # 提供文件下载
        return send_file(file_path, as_attachment=True, download_name=safe_filename)
    except Exception as e:
        logger.error(f"下载文件时出错: {e}")
        return jsonify({'error': '下载文件时出错'}), 500

# 错误处理
@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': '文件未找到'}), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({'error': '服务器内部错误'}), 500

if __name__ == '__main__':
    print(f"启动Flask应用服务器，支持PPTX: {PPTX_AVAILABLE}, 支持PDF: {PDF_AVAILABLE}")
    if not PPTX_AVAILABLE and not PDF_AVAILABLE:
        print("警告: 没有安装python-pptx和reportlab模块，无法生成PPT或PDF文件")
        print("请运行 'pip install python-pptx reportlab' 来安装必要的依赖")
        print("系统将生成文本文件作为替代方案")
    # 使用配置中的host和port
    print(f"服务器配置: host={config['server']['host']}, port={config['server']['port']}")
    app.run(host=config['server']['host'], port=config['server']['port'], debug=False)